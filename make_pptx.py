"""Generate Gemini Workforce Dashboard presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BLUE    = RGBColor(0x25, 0x63, 0xEB)
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
SKY     = RGBColor(0x0E, 0xA5, 0xE9)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
RED     = RGBColor(0xDC, 0x26, 0x26)

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SLATE0  = RGBColor(0xF8, 0xFA, 0xFC)
SLATE1  = RGBColor(0xF1, 0xF5, 0xF9)
SLATE2  = RGBColor(0xE2, 0xE8, 0xF0)
SLATE5  = RGBColor(0x64, 0x74, 0x8B)
SLATE7  = RGBColor(0x33, 0x41, 0x55)
SLATE9  = RGBColor(0x0F, 0x17, 0x2A)

BG_BLUE   = RGBColor(0xEF, 0xF6, 0xFF)
BG_VIOLET = RGBColor(0xF5, 0xF3, 0xFF)
BG_GREEN  = RGBColor(0xF0, 0xFD, 0xF4)

# ── Slide size 16:9 ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper primitives ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, size=Pt(12), bold=False, color=SLATE7,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, size=Pt(12), bold=False, color=SLATE7,
             align=PP_ALIGN.LEFT, space_before=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, tag, title, subtitle=""):
    """Standard slide header band."""
    # top accent bar
    add_rect(slide, 0, 0, W, Inches(0.06), fill=BLUE)
    # tag
    add_text_box(slide, Inches(0.6), Inches(0.18), Inches(9), Inches(0.28),
                 tag, size=Pt(9), bold=True, color=BLUE)
    # title
    add_text_box(slide, Inches(0.6), Inches(0.44), Inches(11.8), Inches(0.72),
                 title, size=Pt(22), bold=True, color=SLATE9)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.12), Inches(11.8), Inches(0.32),
                     subtitle, size=Pt(11), color=SLATE5)
    # gradient rule line (approximated as a thin blue rect)
    add_rect(slide, Inches(0.6), Inches(1.45), Inches(9), Inches(0.04), fill=BLUE)


def card_box(slide, x, y, w, h, title="", bullets=None, fill=WHITE,
             title_color=SLATE5, body_size=Pt(10.5)):
    """Draw a card with optional title and bullet list."""
    add_rect(slide, x, y, w, h, fill=fill, line=SLATE2, line_width=Pt(0.5))
    cy = y + Inches(0.18)
    if title:
        add_text_box(slide, x + Inches(0.15), cy, w - Inches(0.3), Inches(0.22),
                     title, size=Pt(8), bold=True, color=title_color)
        cy += Inches(0.24)
    if bullets:
        txb = slide.shapes.add_textbox(x + Inches(0.15), cy,
                                       w - Inches(0.3), h - (cy - y) - Inches(0.1))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = ("• " if not b.startswith("✓") else "") + b
            run.font.size = body_size
            run.font.color.rgb = SLATE7
            run.font.name = "Calibri"
    return cy


def badge_text(slide, x, y, text, bg, fg):
    bw = Inches(len(text) * 0.085 + 0.25)
    bh = Inches(0.23)
    add_rect(slide, x, y, bw, bh, fill=bg, line=None)
    add_text_box(slide, x + Inches(0.08), y + Inches(0.03), bw - Inches(0.1), bh,
                 text, size=Pt(8), bold=True, color=fg)
    return bw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_BLUE)

# Right panel
add_rect(sl, Inches(9.1), 0, Inches(4.23), H, fill=RGBColor(0xDB, 0xEA, 0xFE))

# Left content
add_text_box(sl, Inches(0.65), Inches(0.5), Inches(6), Inches(0.28),
             "EXECUTIVE BRIEFING · MAY 2026",
             size=Pt(9), bold=True, color=BLUE)

txb = sl.shapes.add_textbox(Inches(0.65), Inches(0.9), Inches(8.2), Inches(2.0))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.add_run().text = "Gemini "
p.runs[0].font.size = Pt(40)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = SLATE9
p.runs[0].font.name = "Calibri"
r2 = p.add_run()
r2.text = "Workforce"
r2.font.size = Pt(40)
r2.font.bold = True
r2.font.color.rgb = BLUE
r2.font.name = "Calibri"

add_para(tf, "Dashboard", size=Pt(40), bold=True, color=SLATE9)

add_text_box(sl, Inches(0.65), Inches(3.0), Inches(8.0), Inches(0.7),
             "AI-powered workforce & spend intelligence —\nnatural language to live charts in seconds, not weeks.",
             size=Pt(13), color=SLATE5)

# Badges row
bx = Inches(0.65)
by = Inches(3.85)
bgs = [(RGBColor(0xDB,0xEA,0xFE), BLUE, "🤖  Gemini 2.5 Flash"),
       (RGBColor(0xE0,0xF2,0xFE), SKY,  "☁️  Cloud Run"),
       (RGBColor(0xED,0xE9,0xFE), VIOLET,"📊  BigQuery"),
       (RGBColor(0xDC,0xFC,0xE7), GREEN, "⚡  Real-time Agent")]
for bg_c, fg_c, txt in bgs:
    bw = badge_text(sl, bx, by, txt, bg_c, fg_c)
    bx += bw + Inches(0.1)

# Meta
add_text_box(sl, Inches(0.65), Inches(4.3), Inches(8), Inches(0.6),
             "Audience: Senior Managers · Architects          Status: Active Development",
             size=Pt(10), color=SLATE5)

# Right panel items
items = [
    ("🤖", "Gemini 2.5 Flash",    "AI Agent · NL→SQL · Chat"),
    ("📊", "Google BigQuery",      "Petabyte-scale Data Warehouse"),
    ("⚛️", "React + TypeScript",   "Drag-drop Dashboard SPA"),
    ("🐍", "FastAPI (Python)",     "Async REST + SSE Backend"),
    ("🐳", "Cloud Run",            "Serverless · Scale-to-Zero"),
]
ry = Inches(0.55)
for icon, name, role in items:
    add_rect(sl, Inches(9.3), ry, Inches(3.8), Inches(1.12),
             fill=WHITE, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, Inches(9.35), ry + Inches(0.08), Inches(0.4), Inches(0.35),
                 icon, size=Pt(18))
    add_text_box(sl, Inches(9.78), ry + Inches(0.06), Inches(3.2), Inches(0.28),
                 name, size=Pt(11), bold=True, color=SLATE7)
    add_text_box(sl, Inches(9.78), ry + Inches(0.34), Inches(3.2), Inches(0.28),
                 role, size=Pt(9), color=SLATE5)
    ry += Inches(1.2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "01 · PROBLEM STATEMENT",
           "Workforce insights are too slow, too manual, too siloed",
           "Every data question becomes a ticket. Every chart takes days. Every manager is blocked.")

cards = [
    ("📋  Reporting Bottleneck",
     ["Ad-hoc spend questions require SQL engineers — 2–5 day turnaround",
      "Monthly workforce reports assembled manually from BigQuery exports",
      "Business leaders make decisions on stale snapshots, not live data"]),
    ("📦  Data Siloed & Opaque",
     ["Workforce, vendor, spend & FTE data spans 3 environments (DEV/UAT/PRD)",
      "No single view — managers switch between BigQuery, Excel, email threads",
      "Schema knowledge locked inside a few engineers' heads"]),
    ("🔐  Access & Governance Friction",
     ["Tableau / PowerBI licenses are expensive and require analyst training",
      "Granting direct BigQuery access to non-technical users is a security risk",
      "No per-user audit trail for data access"]),
    ("🔄  No Self-Service Analytics",
     ['"Top 10 vendors by offshore spend this quarter" — requires a Jira ticket',
      "Dashboard requests pile up in the BI team backlog",
      "Non-technical PMs cannot explore data independently"]),
]
cx, cy_start = Inches(0.35), Inches(1.62)
cw, ch = Inches(6.3), Inches(2.1)
for i, (t, b) in enumerate(cards):
    x = cx + (i % 2) * Inches(6.5)
    y = cy_start + (i // 2) * Inches(2.2)
    card_box(sl, x, y, cw, ch, title=t, bullets=b)

# Warning box
add_rect(sl, Inches(0.35), Inches(6.05), Inches(12.6), Inches(1.0),
         fill=RGBColor(0xFF,0xFB,0xEB), line=RGBColor(0xFD,0xE6,0x8A), line_width=Pt(0.5))
add_text_box(sl, Inches(0.55), Inches(6.12), Inches(0.4), Inches(0.5), "⚠️", size=Pt(18))
add_text_box(sl, Inches(1.0), Inches(6.1), Inches(11.7), Inches(0.28),
             "Business Impact", size=Pt(11), bold=True, color=AMBER)
add_text_box(sl, Inches(1.0), Inches(6.38), Inches(11.7), Inches(0.5),
             "Delayed workforce decisions cost time and budget. When capital vs. expense classifications are wrong by "
             "even one month, it affects financial reporting across the entire portfolio.",
             size=Pt(10), color=RGBColor(0x78,0x35,0x0F))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "02 · SOLUTION",
           "Ask any question in plain English — get a live chart in seconds",
           "Gemini Workforce Dashboard puts real-time BigQuery analytics in the hands of every manager, with zero SQL.")

# Feature highlights row
feats = [
    ("💬", "Natural Language", "Type plain English — AI generates BigQuery SQL & renders a chart"),
    ("🤖", "Conversational AI", "Multi-turn chat with live data — follow-up, refine, drill down"),
    ("📊", "Live Dashboard",   "Drag-drop widgets from chat, pin favourites, save layouts"),
    ("🌐", "Multi-Environment","DEV / UAT / PRD tabs with isolated data & scorecard views"),
    ("📄", "PDF Export",       "One-click branded PDF with AI-generated narrative per widget"),
]
fw = Inches(2.42)
fh = Inches(1.45)
fx = Inches(0.35)
for icon, title, desc in feats:
    add_rect(sl, fx, Inches(1.62), fw, fh, fill=WHITE, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, fx + Inches(0.12), Inches(1.72), fw - Inches(0.2), Inches(0.38),
                 icon, size=Pt(22))
    add_text_box(sl, fx + Inches(0.12), Inches(2.1), fw - Inches(0.2), Inches(0.26),
                 title, size=Pt(10), bold=True, color=SLATE9)
    add_text_box(sl, fx + Inches(0.12), Inches(2.35), fw - Inches(0.2), Inches(0.6),
                 desc, size=Pt(9), color=SLATE5)
    fx += fw + Inches(0.1)

# Two cards below
card_box(sl, Inches(0.35), Inches(3.2), Inches(6.1), Inches(3.9),
         title="Built-in Scorecard Views",
         bullets=["✓ FTE & Hierarchy — headcount, capital vs. expense, monthly FTP",
                  "✓ Vendor Summary — offshore/onshore split, tier breakdown, YTD spend",
                  "✓ Hierarchy Drill — org-level spend waterfall with vendor mix",
                  "✓ Airflow DAG monitor — live pipeline health per environment",
                  "✓ Schema Audit — cross-environment data quality checks",
                  "✓ Excel Mapping — source-to-BigQuery field lineage preview"],
         body_size=Pt(10.5))

card_box(sl, Inches(6.6), Inches(3.2), Inches(6.35), Inches(3.9),
         title="Who Benefits",
         bullets=["👔  Senior Managers — Self-service spend insight without waiting on BI",
                  "🏗️  Architects — Live environment health, schema validation, pipeline monitoring",
                  "📋  Program Managers — Portfolio spend, intake status, resource allocation at a glance",
                  "💰  Finance Teams — Capital vs. expense classification, YTD vs. budget variance"],
         body_size=Pt(10.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "03 · ARCHITECTURE",
           "Single-container deployment on Cloud Run",
           "React SPA + FastAPI backend ship as one Docker image. Scale-to-zero, no idle cost.")

cols = [
    (Inches(0.3),  Inches(4.1),  BG_BLUE,   BLUE,   "BROWSER (CLIENT)",
     ["React SPA (TypeScript + Vite)", "Dashboard Grid — React Grid Layout",
      "Scorecard Tabs: FTE / Vendor / Hierarchy", "AI Chat Panel — SSE streaming",
      "Charts: Recharts · 8 chart types", "DataTable with sort & pagination",
      "Query Bar, Glossary, Favorites"]),
    (Inches(4.55), Inches(4.55), BG_VIOLET, VIOLET, "CLOUD RUN CONTAINER",
     ["FastAPI + uvicorn (Python 3.12)", "POST /api/chat/stream — SSE agent",
      "POST /api/query — NL → SQL → chart", "GET /api/scorecard/{fte,vendor,hierarchy}",
      "POST /api/pdf — Playwright PDF export", "CRUD /api/glossary + /api/favorites",
      "SQLite — glossary & favorites state"]),
    (Inches(9.25), Inches(3.75), BG_GREEN,  TEAL,   "GOOGLE CLOUD SERVICES",
     ["Vertex AI — Gemini 2.5 Flash", "• Function calling / tool use",
      "• Multi-turn chat sessions", "BigQuery — data warehouse",
      "• Schema introspection at startup", "• Ad-hoc query execution",
      "Google OAuth2 — user identity"]),
]

for x, w, bg, col, label, items in cols:
    add_rect(sl, x, Inches(1.62), w, Inches(5.45), fill=bg, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, x + Inches(0.12), Inches(1.72), w - Inches(0.2), Inches(0.22),
                 label, size=Pt(8), bold=True, color=col)
    card_y = Inches(2.0)
    add_rect(sl, x + Inches(0.12), card_y, w - Inches(0.24), Inches(4.8),
             fill=WHITE, line=SLATE2, line_width=Pt(0.3))
    txb = sl.shapes.add_textbox(x + Inches(0.22), card_y + Inches(0.1),
                                 w - Inches(0.44), Inches(4.6))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(9.5)
        r.font.color.rgb = SLATE7
        r.font.name = "Calibri"
        r.font.bold = item.startswith("•") is False and ":" not in item and i == 0

# Arrows
arrow_y = Inches(4.35)
add_text_box(sl, Inches(4.37), arrow_y, Inches(0.3), Inches(0.35), "→", size=Pt(20), color=SLATE5)
add_text_box(sl, Inches(9.05), arrow_y, Inches(0.3), Inches(0.35), "→", size=Pt(20), color=SLATE5)

# Bottom info bar
add_rect(sl, Inches(0.3), Inches(7.15), Inches(12.65), Inches(0.28),
         fill=WHITE, line=SLATE2, line_width=Pt(0.3))
for i, (icon, lbl) in enumerate([("🐳", "One Docker image — React build into FastAPI static/"),
                                   ("📡", "SSE streaming for real-time agent status events"),
                                   ("📐", "Widget layouts in localStorage — no server-side layout DB")]):
    x = Inches(0.45) + i * Inches(4.2)
    add_text_box(sl, x, Inches(7.17), Inches(4.0), Inches(0.24),
                 f"{icon}  {lbl}", size=Pt(8.5), color=SLATE7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AGENT INTERACTION FLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "04 · AGENT INTERACTION",
           "How the AI agent answers a question end-to-end",
           "Two paths: fast conversational (~1s) and an agent loop with live BigQuery tool calls (~3–6s).")

# Fast path column
add_rect(sl, Inches(0.3), Inches(1.62), Inches(6.1), Inches(0.32),
         fill=RGBColor(0xDC,0xFC,0xE7), line=RGBColor(0x86,0xEF,0xAC), line_width=Pt(0.5))
add_text_box(sl, Inches(0.42), Inches(1.66), Inches(6.0), Inches(0.24),
             "⚡  FAST PATH  —  ~1 second  ·  No BigQuery  ·  Conversational questions",
             size=Pt(9), bold=True, color=GREEN)

fast_steps = [
    ("1", "User message arrives",
     "e.g. \"What does YTD mean?\" / \"Explain capital vs. expense\""),
    ("2", "Intent Classifier (0ms)",
     "Regex check — no analytical keywords detected → route to fast path"),
    ("3", "Single Gemini call — chat_turn()",
     "No tools attached, schema context + glossary injected, history preserved"),
    ("4", "JSON → SSE result",
     "{ text, intent:\"explain\", suggested_questions } streamed immediately"),
]
sy = Inches(2.02)
for num, title, desc in fast_steps:
    add_rect(sl, Inches(0.3), sy, Inches(0.3), Inches(0.3), fill=BLUE)
    add_text_box(sl, Inches(0.3), sy + Inches(0.05), Inches(0.3), Inches(0.22),
                 num, size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.65), sy, Inches(5.72), Inches(0.72),
             fill=WHITE, line=SLATE2, line_width=Pt(0.3))
    add_text_box(sl, Inches(0.78), sy + Inches(0.05), Inches(5.5), Inches(0.25),
                 title, size=Pt(10), bold=True, color=SLATE9)
    add_text_box(sl, Inches(0.78), sy + Inches(0.3), Inches(5.5), Inches(0.35),
                 desc, size=Pt(9), color=SLATE5)
    sy += Inches(0.82)

# Agent path column
add_rect(sl, Inches(6.6), Inches(1.62), Inches(6.35), Inches(0.32),
         fill=RGBColor(0xDB,0xEA,0xFE), line=RGBColor(0x93,0xC5,0xFD), line_width=Pt(0.5))
add_text_box(sl, Inches(6.72), Inches(1.66), Inches(6.2), Inches(0.24),
             "🤖  AGENT PATH  —  3–6 seconds  ·  Live BigQuery queries  ·  Data & chart questions",
             size=Pt(9), bold=True, color=BLUE)

agent_steps = [
    ("1", "Intent Classifier → Agent",
     "e.g. \"Show me top 10 vendors by YTD spend\" → analytical keywords detected"),
    ("2", "Contextual Hints (parallel)",
     "Matching keyword columns fetched concurrently via ThreadPoolExecutor — cached 10 min"),
    ("3", "Agent Loop (max 4 rounds)",
     "Gemini reasons → calls run_bigquery_query tool → BQ executes → result back to Gemini"),
    ("4", "Status events stream to UI",
     "\"Analyzing…\" → \"Running 2 queries in parallel (step 1)\" → \"Formatting…\""),
    ("5", "Final JSON + widget data",
     "{ text, intent, widget{sql, chart_type, data}, suggested_questions } — BQ cached"),
]
sy = Inches(2.02)
for num, title, desc in agent_steps:
    add_rect(sl, Inches(6.6), sy, Inches(0.3), Inches(0.3), fill=BLUE)
    add_text_box(sl, Inches(6.6), sy + Inches(0.05), Inches(0.3), Inches(0.22),
                 num, size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(6.95), sy, Inches(5.97), Inches(0.72),
             fill=WHITE, line=SLATE2, line_width=Pt(0.3))
    add_text_box(sl, Inches(7.08), sy + Inches(0.05), Inches(5.75), Inches(0.25),
                 title, size=Pt(10), bold=True, color=SLATE9)
    add_text_box(sl, Inches(7.08), sy + Inches(0.3), Inches(5.75), Inches(0.35),
                 desc, size=Pt(9), color=SLATE5)
    sy += Inches(0.82)

# Tool call box
add_rect(sl, Inches(0.3), Inches(6.3), Inches(12.65), Inches(1.05),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.5), Inches(6.38), Inches(4), Inches(0.22),
             "AGENT TOOL: run_bigquery_query(sql)", size=Pt(8), bold=True, color=SLATE5)
add_text_box(sl, Inches(0.5), Inches(6.6), Inches(9), Inches(0.65),
             "The only tool the agent has is run_bigquery_query(sql). Gemini calls it multiple times per turn — "
             "first SELECT DISTINCT to discover filter values, then to fetch actual data. "
             "All tool calls within a round run in parallel via asyncio.gather.",
             size=Pt(9.5), color=SLATE7)
bx = Inches(9.6)
for txt, bg, fg in [("Tool calls parallelized", RGBColor(0xDB,0xEA,0xFE), BLUE),
                     ("60s result cache",        RGBColor(0xDC,0xFC,0xE7), GREEN),
                     ("Max 4 BQ round-trips",    RGBColor(0xED,0xE9,0xFE), VIOLET)]:
    badge_text(sl, bx, Inches(6.62), txt, bg, fg)
    bx = Inches(9.6)
    if txt == "Tool calls parallelized":
        bx = Inches(9.6)
# stack them vertically instead
bx_items = [("Tool calls parallelized", RGBColor(0xDB,0xEA,0xFE), BLUE),
             ("60s result cache",        RGBColor(0xDC,0xFC,0xE7), GREEN),
             ("Max 4 BQ round-trips",    RGBColor(0xED,0xE9,0xFE), VIOLET)]
by2 = Inches(6.58)
for txt, bg, fg in bx_items:
    badge_text(sl, Inches(9.6), by2, txt, bg, fg)
    by2 += Inches(0.28)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "05 · KEY FEATURES",
           "A full analytics platform — not just a chat interface")

feat_cards = [
    ("🗨️  AI Chat + Streaming",
     ["Multi-turn conversation with full history context",
      "Status events stream in real-time via Server-Sent Events",
      "3 suggested follow-up questions per response",
      "Intent routing: conversational vs. analytical (0ms)",
      "Glossary-aware — custom term→SQL filter mappings"]),
    ("📊  Dynamic Widget Dashboard",
     ["Drag-and-drop layout via React Grid Layout",
      "8 chart types: bar, stacked, line, combo, donut, pie, table, KPI",
      "Add widget directly from chat response",
      "Live / static toggle per widget",
      "Layout persisted per-user in localStorage"]),
    ("📋  Scorecard Views",
     ["FTE & Hierarchy: KPI cards, monthly FTP/spend, capital vs. expense",
      "Vendor Summary: tier breakdown, offshore/onshore, resource count",
      "Hierarchy Drill: org-level vendor spend cascade",
      "5-minute server-side TTL cache — fast repeated loads"]),
    ("🔍  NL Query Bar",
     ["One-line natural language → instant chart widget",
      "Refine existing widget with follow-up instruction",
      "Direct SQL mode for power users",
      "Glossary terms automatically injected into every prompt"]),
    ("📄  PDF Export",
     ["Full dashboard snapshot via Playwright headless Chrome",
      "Branded header with company logo + EST timestamp",
      "AI-generated narrative per widget section",
      "Page breaks, headers & footers on every page"]),
    ("🛠️  Governance Tools",
     ["Glossary manager — map business terms to BQ column filters",
      "Hidden entity mappings (JSON config) for exact SQL injection",
      "Schema Audit — cross-environment column validation",
      "Excel Mapping — source-to-BigQuery field lineage preview",
      "Airflow DAG monitor — pipeline health per environment"]),
]
cw, ch = Inches(4.0), Inches(2.42)
for i, (t, b) in enumerate(feat_cards):
    x = Inches(0.35) + (i % 3) * Inches(4.3)
    y = Inches(1.65) + (i // 3) * Inches(2.58)
    card_box(sl, x, y, cw, ch, title=t, bullets=b, body_size=Pt(10))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "06 · TECHNOLOGY STACK",
           "Modern, cloud-native, open-source foundation",
           "Every component is a Google Cloud or widely-adopted OSS library — no proprietary lock-in below GCP.")

sections = [
    ("🖥️  FRONTEND", BLUE,
     ["React 18 — component UI framework",
      "TypeScript — type-safe codebase",
      "Vite — fast dev server & production build",
      "Tailwind CSS — utility-first styling",
      "Recharts — all 8 chart types",
      "React Grid Layout — drag-drop dashboard",
      "Lucide — icon system",
      "TanStack Query — data fetching & caching"]),
    ("🤖  AI / ML", VIOLET,
     ["Gemini 2.5 Flash — primary model (NL→SQL, agent, PDF)",
      "Vertex AI SDK — GenerativeModel + function calling",
      "Tool Declaration — run_bigquery_query schema",
      "System instructions — schema + few-shot + axis rules",
      "Intent routing — regex classifier (0ms, no API call)",
      "SSE streaming — real-time status events to UI"]),
    ("🐍  BACKEND", TEAL,
     ["FastAPI — async REST API + SSE streaming",
      "uvicorn — ASGI server, production-grade",
      "asyncio + ThreadPoolExecutor — parallel BQ calls",
      "SQLAlchemy + SQLite — glossary & favorites",
      "Playwright — headless Chrome for PDF rendering",
      "Pydantic v2 — request/response validation",
      "httpx — async OAuth2 token validation",
      "openpyxl — Excel file parsing"]),
    ("☁️  INFRASTRUCTURE", GREEN,
     ["Cloud Run — serverless container, scale-to-zero",
      "Google BigQuery — petabyte-scale OLAP warehouse",
      "Vertex AI API — managed Gemini endpoints",
      "Google OAuth2 — user identity + BQ credential scoping",
      "Docker — multi-stage image (React → Python)",
      "python-dotenv — environment configuration"]),
]
cw, ch = Inches(5.9), Inches(5.55)
for i, (title, col, items) in enumerate(sections):
    x = Inches(0.35) + (i % 2) * Inches(6.45)
    y = Inches(1.65) + (i // 2) * Inches(2.82)
    add_rect(sl, x, y, cw, ch, fill=WHITE, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, x + Inches(0.14), y + Inches(0.1), cw - Inches(0.25), Inches(0.22),
                 title, size=Pt(8.5), bold=True, color=col)
    txb = sl.shapes.add_textbox(x + Inches(0.14), y + Inches(0.35), cw - Inches(0.25),
                                  ch - Inches(0.45))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = "• " + item
        r.font.size = Pt(10)
        r.font.color.rgb = SLATE7
        r.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — APIS & SERVICES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "07 · APIs & SERVICES",
           "Backend API surface & external service integrations")

# REST API card
add_rect(sl, Inches(0.3), Inches(1.65), Inches(5.9), Inches(5.6),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(1.75), Inches(5.6), Inches(0.22),
             "INTERNAL REST API (FastAPI)", size=Pt(8), bold=True, color=SLATE5)

endpoints = [
    ("POST", BLUE,   "/api/query",                "NL → SQL → chart widget"),
    ("POST", BLUE,   "/api/query/sql",             "Direct SQL execution"),
    ("POST", BLUE,   "/api/chat",                  "AI agent (sync)"),
    ("POST", BLUE,   "/api/chat/stream",            "AI agent (SSE streaming)"),
    ("GET",  GREEN,  "/api/scorecard/fte",          "FTE scorecard data bundle"),
    ("GET",  GREEN,  "/api/scorecard/vendor",        "Vendor scorecard bundle"),
    ("GET",  GREEN,  "/api/scorecard/hierarchy",     "Hierarchy scorecard"),
    ("POST", BLUE,   "/api/pdf",                    "Generate branded PDF report"),
    ("CRUD", SLATE5, "/api/glossary",               "Business term management"),
    ("CRUD", SLATE5, "/api/favorites",              "Saved widget management"),
    ("GET",  GREEN,  "/api/health",                 "AI + BQ readiness check"),
]
ey = Inches(2.02)
for method, col, path, desc in endpoints:
    add_rect(sl, Inches(0.45), ey, Inches(0.55), Inches(0.26), fill=col)
    add_text_box(sl, Inches(0.45), ey + Inches(0.04), Inches(0.55), Inches(0.2),
                 method, size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(1.05), ey + Inches(0.04), Inches(4.65), Inches(0.22),
                 f"{path}  —  {desc}", size=Pt(9.5), color=SLATE7)
    ey += Inches(0.42)

# External services
add_rect(sl, Inches(6.4), Inches(1.65), Inches(6.55), Inches(3.8),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.55), Inches(1.75), Inches(6.3), Inches(0.22),
             "EXTERNAL GCP SERVICES", size=Pt(8), bold=True, color=SLATE5)

ext = [
    ("🤖", "Vertex AI — Gemini 2.5 Flash",
     "Function calling, multi-turn chat, streaming generation.\nConfigured via VERTEX_AI_PROJECT + GEMINI_MODEL env vars."),
    ("📊", "BigQuery API",
     "Ad-hoc query execution, schema introspection.\nBilling project configurable independently from data project."),
    ("🔑", "Google OAuth2 Userinfo API",
     "GET googleapis.com/oauth2/v2/userinfo\nResolves Bearer token → user ID + email for data scoping."),
]
esy = Inches(2.05)
for icon, title, desc in ext:
    add_text_box(sl, Inches(6.55), esy, Inches(0.4), Inches(0.35), icon, size=Pt(16))
    add_text_box(sl, Inches(7.0), esy, Inches(5.7), Inches(0.26),
                 title, size=Pt(10.5), bold=True, color=SLATE9)
    add_text_box(sl, Inches(7.0), esy + Inches(0.28), Inches(5.7), Inches(0.6),
                 desc, size=Pt(9.5), color=SLATE5)
    add_rect(sl, Inches(6.55), esy + Inches(1.05), Inches(6.1), Inches(0.01), fill=SLATE2)
    esy += Inches(1.18)

# Optional integrations
add_rect(sl, Inches(6.4), Inches(5.6), Inches(6.55), Inches(1.65),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.55), Inches(5.7), Inches(6.3), Inches(0.22),
             "OPTIONAL / FUTURE INTEGRATIONS", size=Pt(8), bold=True, color=SLATE5)
future = ["Apache Airflow REST API", "Cloud SQL (replace SQLite)",
          "BigQuery BI Engine", "Cloud Scheduler (reports)", "Secret Manager"]
bx2, by2 = Inches(6.55), Inches(5.98)
for txt in future:
    bw = badge_text(sl, bx2, by2, txt, SLATE1, SLATE7)
    bx2 += bw + Inches(0.08)
    if bx2 > Inches(12.5):
        bx2 = Inches(6.55)
        by2 += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SECURITY & AUTHORIZATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "08 · SECURITY & AUTHORIZATION",
           "Defense-in-depth: user-scoped credentials at every layer",
           "No user ever gets broader BigQuery access than their Google identity already has.")

# Auth flow steps
add_rect(sl, Inches(0.3), Inches(1.65), Inches(6.1), Inches(5.6),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(1.75), Inches(5.8), Inches(0.22),
             "AUTHENTICATION FLOW", size=Pt(8), bold=True, color=SLATE5)

auth_steps = [
    ("1", "Browser sends Bearer token",
     "Google OAuth2 access token in Authorization header on every API request"),
    ("2", "Backend validates token",
     "Calls Google /userinfo endpoint — extracts user ID + email. Falls back to anonymous in dev mode."),
    ("3", "Token forwarded to BigQuery",
     "Same Bearer token wraps into google.oauth2.Credentials — BQ runs queries under user's own GCP IAM identity"),
    ("4", "User-scoped data isolation",
     "Glossary terms and favorites in SQLite are filtered by user ID — no cross-user data leakage"),
]
asy = Inches(2.05)
for num, title, desc in auth_steps:
    add_rect(sl, Inches(0.45), asy, Inches(0.35), Inches(0.35), fill=BLUE)
    add_text_box(sl, Inches(0.45), asy + Inches(0.07), Inches(0.35), Inches(0.22),
                 num, size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.85), asy, Inches(5.4), Inches(1.12),
             fill=SLATE0, line=SLATE1, line_width=Pt(0.3))
    add_text_box(sl, Inches(0.98), asy + Inches(0.06), Inches(5.15), Inches(0.26),
                 title, size=Pt(10.5), bold=True, color=SLATE9)
    add_text_box(sl, Inches(0.98), asy + Inches(0.34), Inches(5.15), Inches(0.65),
                 desc, size=Pt(9.5), color=SLATE5)
    asy += Inches(1.22)

# Security controls
add_rect(sl, Inches(6.6), Inches(1.65), Inches(6.35), Inches(3.8),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(1.75), Inches(6.1), Inches(0.22),
             "SECURITY CONTROLS", size=Pt(8), bold=True, color=SLATE5)

controls = [
    ("🛡️", "CORS Whitelist",
     "Allowed origins explicitly listed: localhost:5173, localhost:3000, CLOUD_RUN_URL. No wildcard (*) in production."),
    ("🔒", "No Direct SQL from Users",
     "All SQL is AI-generated from natural language. Only admins use /api/query/sql. Users cannot construct arbitrary queries."),
    ("📝", "BigQuery Job Audit Trail",
     "Every BQ query runs as a job attributed to the user's GCP identity — full audit in Cloud Logging automatically."),
]
csy = Inches(2.05)
for icon, title, desc in controls:
    add_text_box(sl, Inches(6.75), csy, Inches(0.3), Inches(0.3), icon, size=Pt(14))
    add_text_box(sl, Inches(7.1), csy, Inches(5.65), Inches(0.26),
                 title, size=Pt(10.5), bold=True, color=SLATE9)
    add_text_box(sl, Inches(7.1), csy + Inches(0.28), Inches(5.65), Inches(0.65),
                 desc, size=Pt(9.5), color=SLATE5)
    add_rect(sl, Inches(6.75), csy + Inches(1.05), Inches(5.95), Inches(0.01), fill=SLATE2)
    csy += Inches(1.18)

# Compliance box
add_rect(sl, Inches(6.6), Inches(5.6), Inches(6.35), Inches(1.65),
         fill=BG_GREEN, line=RGBColor(0x86,0xEF,0xAC), line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(5.7), Inches(6.1), Inches(0.22),
             "COMPLIANCE NOTES", size=Pt(8), bold=True, color=GREEN)
comp = ["✓  No PII stored in application DB — only user IDs and term definitions",
        "✓  BigQuery row-level security applied independently at the GCP layer",
        "✓  All secrets via .env / environment variables — never hardcoded"]
cy2 = Inches(5.98)
for c in comp:
    add_text_box(sl, Inches(6.75), cy2, Inches(6.0), Inches(0.26),
                 c, size=Pt(10), color=SLATE7)
    cy2 += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PERFORMANCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "09 · PERFORMANCE",
           "Six layers of optimization — from 8s to sub-2s for common queries")

# Latency table
add_rect(sl, Inches(0.3), Inches(1.65), Inches(6.1), Inches(5.6),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(1.75), Inches(5.8), Inches(0.22),
             "LATENCY IMPROVEMENTS", size=Pt(8), bold=True, color=SLATE5)

# Header row
for x, lbl in [(Inches(0.45), "Optimization"), (Inches(3.75), "Before"),
               (Inches(4.45), "After")]:
    add_text_box(sl, x, Inches(2.0), Inches(1.2), Inches(0.22),
                 lbl, size=Pt(8), bold=True, color=SLATE5)

perf = [
    ("Contextual hint fetching (3 keywords matched)", "~9s seq.", "~3s par.", 0.67),
    ("Repeated hint queries (same keyword)",           "~3s",     "0ms",     1.0),
    ("Widget BQ query (was run twice)",                "2× BQ",   "1× BQ",  0.50),
    ("Agent max tool rounds",                          "10 rnd",  "4 rnd",  0.60),
    ("Conversational questions (no BQ needed)",        "~5s",     "~1s",    0.80),
    ("Scorecard page load (repeat visit)",             "~8s",     "<1s",    0.88),
]
py = Inches(2.28)
for label, before, after, pct in perf:
    add_rect(sl, Inches(0.45), py, Inches(5.4), Inches(0.5), fill=SLATE0, line=SLATE1, line_width=Pt(0.3))
    add_text_box(sl, Inches(0.55), py + Inches(0.06), Inches(3.1), Inches(0.4),
                 label, size=Pt(9), color=SLATE7)
    add_text_box(sl, Inches(3.75), py + Inches(0.12), Inches(0.7), Inches(0.26),
                 before, size=Pt(9), bold=True, color=RED)
    add_text_box(sl, Inches(4.45), py + Inches(0.12), Inches(0.7), Inches(0.26),
                 after, size=Pt(9), bold=True, color=GREEN)
    # mini progress bar
    add_rect(sl, Inches(5.2), py + Inches(0.2), Inches(0.55), Inches(0.1), fill=SLATE2)
    add_rect(sl, Inches(5.2), py + Inches(0.2), Inches(0.55 * pct), Inches(0.1), fill=GREEN)
    py += Inches(0.54)

# Cache architecture
add_rect(sl, Inches(6.6), Inches(1.65), Inches(6.35), Inches(3.2),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(1.75), Inches(6.1), Inches(0.22),
             "CACHING ARCHITECTURE", size=Pt(8), bold=True, color=SLATE5)

caches = [
    ("Schema cache",        "BQ column definitions",         "∞ (process lifetime)",  BLUE),
    ("Hint cache",          "SELECT DISTINCT per column",     "10 minutes",            GREEN),
    ("BQ result cache",     "Agent tool call results",        "60 seconds",            VIOLET),
    ("Scorecard cache",     "Pre-built dashboard queries",    "5 minutes",             AMBER),
    ("Agent model",         "GenerativeModel instance",       "∞ (process lifetime)",  SKY),
]
ccy = Inches(2.02)
for name, desc, ttl, col in caches:
    add_rect(sl, Inches(6.75), ccy, Inches(6.0), Inches(0.5),
             fill=SLATE0, line=SLATE1, line_width=Pt(0.3))
    add_text_box(sl, Inches(6.88), ccy + Inches(0.05), Inches(3.2), Inches(0.22),
                 name, size=Pt(10), bold=True, color=SLATE9)
    add_text_box(sl, Inches(6.88), ccy + Inches(0.27), Inches(3.2), Inches(0.2),
                 desc, size=Pt(8.5), color=SLATE5)
    badge_text(sl, Inches(10.1), ccy + Inches(0.12), ttl,
               RGBColor(0xDB,0xEA,0xFE), col)
    ccy += Inches(0.56)

# Concurrency card
add_rect(sl, Inches(6.6), Inches(5.0), Inches(6.35), Inches(2.25),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(5.1), Inches(6.1), Inches(0.22),
             "CONCURRENCY DESIGN", size=Pt(8), bold=True, color=SLATE5)
conc = ["asyncio event loop — all I/O non-blocking",
        "BQ tool calls within a round run in parallel (asyncio.gather)",
        "Hint queries fire concurrently (ThreadPoolExecutor)",
        "Scorecard endpoints: 6–9 parallel BQ calls via asyncio.gather",
        "Schema pre-warmed in background thread at startup"]
concy = Inches(5.38)
for c in conc:
    add_text_box(sl, Inches(6.75), concy, Inches(6.1), Inches(0.26),
                 "• " + c, size=Pt(10), color=SLATE7)
    concy += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — COST
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "10 · COST OF DEPLOYMENT",
           "Estimated running cost — ~$20–60 / month for a small active team",
           "Cloud Run scales to zero. You pay only for what you use.")

# Cost table
add_rect(sl, Inches(0.3), Inches(1.65), Inches(6.1), Inches(4.0),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(1.75), Inches(5.8), Inches(0.22),
             "MONTHLY COST (est. 10 active users)", size=Pt(8), bold=True, color=SLATE5)

add_rect(sl, Inches(0.45), Inches(2.02), Inches(5.8), Inches(0.28), fill=SLATE1)
for x, lbl in [(Inches(0.55), "Service"), (Inches(2.4), "Usage Assumption"),
               (Inches(4.95), "Est. Cost")]:
    add_text_box(sl, x, Inches(2.05), Inches(2.0), Inches(0.22),
                 lbl, size=Pt(8), bold=True, color=SLATE5)

cost_rows = [
    ("Cloud Run",            "2 vCPU / 2 GB · 8 hrs/day · 22 days",   "$8–15"),
    ("Gemini 2.5 Flash",     "~500 agent calls/mo · ~3K tokens avg",   "$0.50–2"),
    ("BigQuery on-demand",   "~2,000 queries · 500 MB avg scan",        "$5–12"),
    ("Vertex AI API",        "Included in GCP project quota",           "$0"),
    ("Cloud Logging",        "First 50 GB/month free",                  "$0"),
]
cty = Inches(2.35)
for i, (svc, usage, cost) in enumerate(cost_rows):
    bg = SLATE0 if i % 2 == 0 else WHITE
    add_rect(sl, Inches(0.45), cty, Inches(5.8), Inches(0.38), fill=bg)
    add_text_box(sl, Inches(0.55), cty + Inches(0.07), Inches(1.8), Inches(0.26),
                 svc, size=Pt(10), bold=True, color=SLATE7)
    add_text_box(sl, Inches(2.4), cty + Inches(0.07), Inches(2.5), Inches(0.26),
                 usage, size=Pt(9.5), color=SLATE5)
    add_text_box(sl, Inches(4.95), cty + Inches(0.07), Inches(1.2), Inches(0.26),
                 cost, size=Pt(10), bold=True, color=SLATE7)
    cty += Inches(0.38)

# Total row
add_rect(sl, Inches(0.45), cty, Inches(5.8), Inches(0.4), fill=BG_BLUE)
add_text_box(sl, Inches(0.55), cty + Inches(0.08), Inches(1.8), Inches(0.26),
             "TOTAL", size=Pt(10), bold=True, color=BLUE)
add_text_box(sl, Inches(2.4), cty + Inches(0.08), Inches(2.5), Inches(0.26),
             "10 users, moderate daily usage", size=Pt(9.5), color=SLATE5)
add_text_box(sl, Inches(4.95), cty + Inches(0.08), Inches(1.2), Inches(0.26),
             "~$14–29/mo", size=Pt(10), bold=True, color=BLUE)

# Notes box
add_rect(sl, Inches(0.3), Inches(5.82), Inches(6.1), Inches(1.43),
         fill=RGBColor(0xF0,0xF9,0xFF), line=RGBColor(0x7D,0xD3,0xFC), line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(5.92), Inches(5.8), Inches(0.22),
             "COST SCALE NOTES", size=Pt(8), bold=True, color=SKY)
notes = ["Gemini 2.5 Flash is ~10× cheaper than 2.5 Pro — significant savings at scale",
         "Intent routing reduces BQ scans ~30% by skipping non-data questions",
         "Scale-to-zero on Cloud Run means $0 cost outside business hours"]
ny = Inches(6.18)
for n in notes:
    add_text_box(sl, Inches(0.45), ny, Inches(5.8), Inches(0.26),
                 "• " + n, size=Pt(9.5), color=SLATE7)
    ny += Inches(0.3)

# Comparison table
add_rect(sl, Inches(6.6), Inches(1.65), Inches(6.35), Inches(5.6),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(1.75), Inches(6.1), Inches(0.22),
             "COST VS. ALTERNATIVES", size=Pt(8), bold=True, color=SLATE5)

add_rect(sl, Inches(6.75), Inches(2.02), Inches(6.05), Inches(0.28), fill=SLATE1)
for x, lbl in [(Inches(6.82), "Tool"), (Inches(8.8), "Self-Svc"),
               (Inches(9.5), "AI Chat"), (Inches(10.25), "Live BQ"),
               (Inches(11.1), "Monthly Cost")]:
    add_text_box(sl, x, Inches(2.05), Inches(1.3), Inches(0.22),
                 lbl, size=Pt(7.5), bold=True, color=SLATE5)

cmp_rows = [
    ("This Dashboard ✓", "✓", "✓", "✓", "~$20–60",  True),
    ("Tableau (10 users)", "✓", "✗", "~", "$700+",   False),
    ("PowerBI Pro (10 u)", "✓", "✗", "~", "$100+",   False),
    ("Looker (GCP)",       "✓", "~", "✓", "$5K+/yr", False),
    ("Custom BI eng team", "✗", "✗", "✓", "$15K+/mo",False),
]
cry = Inches(2.35)
for row in cmp_rows:
    tool, ss, ai, bq, cost, hl = row
    fill = BG_BLUE if hl else (SLATE0 if cmp_rows.index(row) % 2 == 0 else WHITE)
    add_rect(sl, Inches(6.75), cry, Inches(6.05), Inches(0.45), fill=fill)
    add_text_box(sl, Inches(6.82), cry + Inches(0.09), Inches(1.9), Inches(0.26),
                 tool, size=Pt(10), bold=hl, color=BLUE if hl else SLATE7)
    for x, val in [(Inches(8.88), ss), (Inches(9.6), ai),
                   (Inches(10.35), bq), (Inches(11.1), cost)]:
        col = GREEN if val == "✓" else (RED if val == "✗" else AMBER)
        if "$" in val:
            col = SLATE7 if not hl else BLUE
        add_text_box(sl, x, cry + Inches(0.09), Inches(1.0), Inches(0.26),
                     val, size=Pt(10), bold=(val == "✓" or hl), color=col)
    cry += Inches(0.5)

# Gemini pricing
add_rect(sl, Inches(6.6), Inches(5.62), Inches(6.35), Inches(1.63),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(5.72), Inches(6.1), Inches(0.22),
             "GEMINI 2.5 FLASH PRICING (VERTEX AI)", size=Pt(8), bold=True, color=SLATE5)
pricing = [("Input tokens", "$0.075 / 1M tokens"),
           ("Output tokens", "$0.30 / 1M tokens"),
           ("500 calls/month (~3K tokens avg)", "~$0.50 / month")]
pry = Inches(5.98)
for k, v in pricing:
    add_text_box(sl, Inches(6.75), pry, Inches(3.5), Inches(0.26),
                 k, size=Pt(10), color=SLATE7)
    add_text_box(sl, Inches(10.3), pry, Inches(2.3), Inches(0.26),
                 v, size=Pt(10), bold=True, color=SLATE9)
    pry += Inches(0.3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "11 · DEPLOYMENT",
           "One Docker image, one Cloud Run service, zero managed infrastructure")

# Dockerfile
add_rect(sl, Inches(0.3), Inches(1.65), Inches(6.1), Inches(2.7),
         fill=SLATE9, line=SLATE7, line_width=Pt(0.5))
add_text_box(sl, Inches(0.5), Inches(1.75), Inches(5.7), Inches(0.22),
             "DOCKERFILE (MULTI-STAGE)", size=Pt(8), bold=True,
             color=RGBColor(0x94,0xA3,0xB8))
dockerfile = (
    "# Stage 1: Build React SPA\n"
    "FROM node:20 AS ui-build\n"
    "WORKDIR /app/frontend\n"
    "RUN npm ci && npm run build\n\n"
    "# Stage 2: Python backend\n"
    "FROM python:3.12-slim\n"
    "COPY requirements.txt .\n"
    "RUN pip install -r requirements.txt\n"
    "COPY --from=ui-build /app/frontend/dist ./static\n"
    'CMD ["uvicorn", "main:app", "--host", "0.0.0.0", "--port", "8080"]'
)
txb = sl.shapes.add_textbox(Inches(0.5), Inches(2.02), Inches(5.7), Inches(2.2))
txb.word_wrap = False
tf = txb.text_frame
tf.word_wrap = False
for i, line in enumerate(dockerfile.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.size = Pt(8.5)
    r.font.name = "Courier New"
    r.font.color.rgb = RGBColor(0x86,0xEF,0xAC) if line.startswith("#") else RGBColor(0xE2,0xE8,0xF0)

# Env vars
add_rect(sl, Inches(0.3), Inches(4.5), Inches(6.1), Inches(2.75),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(0.45), Inches(4.6), Inches(5.8), Inches(0.22),
             "CONFIGURATION (ENVIRONMENT VARIABLES)", size=Pt(8), bold=True, color=SLATE5)
envs = [("VERTEX_AI_PROJECT",   "my-gcp-project"),
        ("VERTEX_AI_LOCATION",  "us-central1"),
        ("GEMINI_MODEL",        "gemini-2.5-flash"),
        ("BIGQUERY_PROJECT_ID", "my-data-project"),
        ("BIGQUERY_TABLES",     "proj.dataset.table1,..."),
        ("CLOUD_RUN_URL",       "https://app-xxx.run.app")]
evy = Inches(4.88)
for k, v in envs:
    add_text_box(sl, Inches(0.45), evy, Inches(2.3), Inches(0.26),
                 k, size=Pt(9.5), bold=True, color=BLUE)
    add_text_box(sl, Inches(2.8), evy, Inches(3.4), Inches(0.26),
                 "= " + v, size=Pt(9.5), color=SLATE7)
    evy += Inches(0.3)

# Cloud Run config
add_rect(sl, Inches(6.6), Inches(1.65), Inches(6.35), Inches(3.1),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(1.75), Inches(6.1), Inches(0.22),
             "CLOUD RUN CONFIGURATION", size=Pt(8), bold=True, color=SLATE5)
cr_config = [("Container port", "8080"),
             ("Min instances",  "0 (scale-to-zero)"),
             ("Max instances",  "3–5 (configurable)"),
             ("CPU",            "2 vCPU (always-on for SSE)"),
             ("Memory",         "2 GB"),
             ("Concurrency",    "80 requests / instance"),
             ("Request timeout","300s (PDF + long agent calls)")]
ccy2 = Inches(2.02)
for k, v in cr_config:
    add_rect(sl, Inches(6.75), ccy2, Inches(6.05), Inches(0.36), fill=SLATE0)
    add_text_box(sl, Inches(6.85), ccy2 + Inches(0.06), Inches(3.0), Inches(0.26),
                 k, size=Pt(10), color=SLATE7)
    add_text_box(sl, Inches(9.9), ccy2 + Inches(0.06), Inches(2.7), Inches(0.26),
                 v, size=Pt(10), bold=True, color=SLATE9)
    ccy2 += Inches(0.38)

# IAM roles
add_rect(sl, Inches(6.6), Inches(4.9), Inches(6.35), Inches(2.35),
         fill=WHITE, line=SLATE2, line_width=Pt(0.5))
add_text_box(sl, Inches(6.75), Inches(5.0), Inches(6.1), Inches(0.22),
             "IAM ROLES REQUIRED", size=Pt(8), bold=True, color=SLATE5)
iam = [("Cloud Run SA", VIOLET,
        "roles/bigquery.jobUser + roles/aiplatform.user"),
       ("End Users",    BLUE,
        "roles/bigquery.dataViewer (on the data tables)"),
       ("Vertex AI",    GREEN,
        "roles/aiplatform.user on VERTEX_AI_PROJECT")]
iamy = Inches(5.28)
for entity, col, role in iam:
    badge_text(sl, Inches(6.75), iamy, entity, SLATE1, col)
    add_text_box(sl, Inches(7.7), iamy, Inches(5.0), Inches(0.26),
                 role, size=Pt(9.5), color=SLATE7)
    iamy += Inches(0.38)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, SLATE0)
add_header(sl, "12 · ROADMAP", "What's next — planned enhancements")

near = [
    ("BigQuery BI Engine",
     "Reserve in-memory capacity → query latency drops from 2–4s to 50–200ms for the main dataset"),
    ("Materialized Views for Scorecards",
     "Pre-aggregate heavy GROUP BY queries so scorecard data loads in milliseconds, not seconds"),
    ("Airflow DAG Integration",
     "Backend proxy to Airflow REST API — DAG graph, run history, and SQL viewer in the dashboard"),
    ("Gemini Response Streaming",
     "Stream partial narrative text to the UI while BQ queries run — perceived latency near zero"),
]
far = [
    ("Role-Based Access Control",
     "Manager / analyst / admin roles controlling which scorecard tabs and data subsets are visible"),
    ("Scheduled Report Delivery",
     "Cloud Scheduler + SendGrid — weekly PDF scorecards auto-emailed to stakeholders"),
    ("Multi-Table Query Support",
     "JOIN across workforce + financial + intake tables — agent selects tables based on question context"),
    ("Cloud SQL (PostgreSQL)",
     "Replace SQLite with managed Cloud SQL for multi-instance deployments and shared glossary"),
]

for col_items, col_x, col_title, dot_col in [
    (near, Inches(0.3), "NEAR-TERM  (Next 90 Days)", BLUE),
    (far,  Inches(6.6), "MEDIUM-TERM  (90–180 Days)", VIOLET),
]:
    add_rect(sl, col_x, Inches(1.65), Inches(6.1), Inches(5.3),
             fill=WHITE, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, col_x + Inches(0.15), Inches(1.75), Inches(5.8), Inches(0.22),
                 col_title, size=Pt(8), bold=True, color=dot_col)
    tl_y = Inches(2.1)
    for title, desc in col_items:
        add_rect(sl, col_x + Inches(0.18), tl_y + Inches(0.08),
                 Inches(0.14), Inches(0.14), fill=dot_col)
        add_rect(sl, col_x + Inches(0.245), tl_y + Inches(0.22),
                 Inches(0.01), Inches(0.9), fill=SLATE2)
        add_text_box(sl, col_x + Inches(0.42), tl_y, Inches(5.5), Inches(0.26),
                     title, size=Pt(10.5), bold=True, color=SLATE9)
        add_text_box(sl, col_x + Inches(0.42), tl_y + Inches(0.28), Inches(5.5), Inches(0.55),
                     desc, size=Pt(9.5), color=SLATE5)
        tl_y += Inches(1.12)

# Vision bar
add_rect(sl, Inches(0.3), Inches(7.08), Inches(12.65), Inches(0.32),
         fill=BG_VIOLET, line=RGBColor(0xC4,0xB5,0xFD), line_width=Pt(0.5))
add_text_box(sl, Inches(0.5), Inches(7.11), Inches(0.3), Inches(0.26), "🚀", size=Pt(13))
add_text_box(sl, Inches(0.85), Inches(7.12), Inches(11.9), Inches(0.24),
             "Vision: Any Google-authenticated user asks workforce data questions in plain English — "
             "answers in under 2 seconds, governed by existing BigQuery IAM roles.",
             size=Pt(9.5), color=SLATE7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, BG_BLUE)
add_rect(sl, 0, 0, W, Inches(0.08), fill=BLUE)

add_text_box(sl, Inches(1.5), Inches(0.7), Inches(10), Inches(0.28),
             "THANK YOU", size=Pt(10), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

txb = sl.shapes.add_textbox(Inches(1.5), Inches(1.05), Inches(10), Inches(1.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Ask your data anything."
r.font.size = Pt(34)
r.font.bold = True
r.font.color.rgb = SLATE9
r.font.name = "Calibri"

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Get answers in seconds."
r2.font.size = Pt(34)
r2.font.bold = True
r2.font.color.rgb = BLUE
r2.font.name = "Calibri"

add_text_box(sl, Inches(2.5), Inches(2.72), Inches(8.33), Inches(0.62),
             "Gemini Workforce Dashboard puts the full power of BigQuery + Gemini in the "
             "hands of every manager — no SQL, no tickets, no waiting.",
             size=Pt(13), color=SLATE5, align=PP_ALIGN.CENTER)

# Stat boxes
stats = [("~1s", "Fast path\nresponse", BLUE),
         ("~3–6s", "Agent with\nlive BQ data", GREEN),
         ("$20–60", "Est. monthly\ncost (10 users)", VIOLET),
         ("0", "SQL needed\nby end users", TEAL)]
bx = Inches(1.0)
for num, lbl, col in stats:
    add_rect(sl, bx, Inches(3.5), Inches(2.6), Inches(1.35),
             fill=WHITE, line=SLATE2, line_width=Pt(0.5))
    add_text_box(sl, bx, Inches(3.6), Inches(2.6), Inches(0.55),
                 num, size=Pt(28), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(sl, bx + Inches(0.1), Inches(4.15), Inches(2.4), Inches(0.6),
                 lbl, size=Pt(10), color=SLATE5, align=PP_ALIGN.CENTER)
    bx += Inches(2.83)

# Tech badges
tech = [("React + TypeScript", RGBColor(0xDB,0xEA,0xFE), BLUE),
        ("FastAPI + Python",   RGBColor(0xED,0xE9,0xFE), VIOLET),
        ("Gemini 2.5 Flash",   RGBColor(0xDC,0xFC,0xE7), GREEN),
        ("Google BigQuery",    RGBColor(0xE0,0xF2,0xFE), SKY),
        ("Cloud Run",          RGBColor(0xFE,0xF3,0xC7), AMBER)]
total_w = sum(Inches(len(t[0]) * 0.088 + 0.3) for t in tech) + Inches(0.12 * (len(tech) - 1))
bx = (W - total_w) / 2
by3 = Inches(5.1)
for txt, bg, fg in tech:
    bw = badge_text(sl, bx, by3, txt, bg, fg)
    bx += bw + Inches(0.12)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/navin/gemini-dashboard/GeminiWorkforceDashboard.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"    Slides: {len(prs.slides)}")
